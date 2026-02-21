"""
PPTX Loader — one Document per slide.

Why slide-per-document:
- A slide is a natural semantic unit (one idea, one topic)
- Title + bullets together have more meaning than either alone
- Slide number is useful metadata ("see slide 7 for pricing breakdown")

What we extract per slide:
- Title (if present)
- All text from shapes (bullets, text boxes, notes)
- Tables within slides (serialized row-by-row)
- Slide number

Speaker notes: extracted and appended — they often contain richer explanation
than the slide bullets themselves.
"""
from pathlib import Path

from langchain_core.documents import Document

from app.loaders.base import BaseLoader


class PptxLoader(BaseLoader):
    SUPPORTED_EXTENSIONS = [".pptx", ".ppt"]

    def load(self, file_path: Path) -> list[Document]:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        docs = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            parts = []
            title = ""
            notes_text = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.TITLE
                        title = text
                    elif text:
                        parts.append(text)

                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        rows.append(" | ".join(
                            cell.text_frame.text.strip() for cell in row.cells
                        ))
                    parts.append("TABLE:\n" + "\n".join(rows))

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame
                if notes:
                    notes_text = notes.text.strip()

            content_parts = []
            if title:
                content_parts.append(f"Slide {slide_num}: {title}")
            content_parts.extend(parts)
            if notes_text:
                content_parts.append(f"Notes: {notes_text}")

            content = "\n".join(content_parts)
            if content.strip():
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": file_path.name,
                        "slide": slide_num,
                        "title": title,
                    },
                ))

        return docs
